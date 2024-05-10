from fastapi import FastAPI, HTTPException, Request, Response
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
import requests
import base64
import bs4
import dataclasses
import threading
import time
import pptx
import io
import uvicorn
import PIL
from fp.fp import FreeProxy
import os

USE_PROXIES = os.environ["USE_PROXY"] if "USE_PROXY" in os.environ else False


def decode_b64(message: str) -> str:
    base64_message = message
    base64_bytes = base64_message.encode("ascii")
    message_bytes = base64.b64decode(base64_bytes)
    message = message_bytes.decode("ascii")

    return message


@dataclasses.dataclass
class ImageSrcSet:
    px_size: int = -1
    url: str = ""


def scrape_slides_number(soup: bs4.BeautifulSoup) -> int:

    page_number_spans = soup.find_all("span", {"data-cy": "page-number"})
    if not page_number_spans:
        return -1

    page_number_span_text = page_number_spans[0].text

    if " of " not in page_number_span_text:
        return -1

    cur_slide_num, total_slide_num = page_number_span_text.split(" of ")

    return int(total_slide_num)


def scrape_highest_img_src_set(soup: bs4.BeautifulSoup) -> ImageSrcSet:

    ret = ImageSrcSet()

    first_slide_tags = soup.find_all("img", {"id": "slide-image-0"})
    if not first_slide_tags:
        print("'img' tag with 'id' attribute set to 'slide-image-0' not found")
        return ret

    try:
        first_slide_tag = first_slide_tags[0]["srcset"]
    except:
        print(
            "Could not get 'srcset' attribute of 'img' tag with 'id' attribute set to 'slide-image-0'"
        )
        return ret

    if "," not in first_slide_tag:
        print(first_slide_tag)
        return ret

    qualities_list = first_slide_tag.split(",")

    if not qualities_list:
        print("No available qualities")
        return ret

    highest_quality = qualities_list[-1]
    _, url, px_size = highest_quality.split(" ")

    ret.px_size = int(px_size.replace("w", ""))

    url_1, url2 = url.split(f"-1-{ret.px_size}")
    ret.url = url_1 + f"-SLIDE_NUMBER-{ret.px_size}" + url2

    return ret


def save_image(url: str, idx: int, downloaded_images: list, lock, proxy=None):

    proxies = {}
    if USE_PROXIES:
        proxies = {"http": proxy}

    req = requests.get(
        url,
        stream=True,
        proxies=proxies,
    )
    req.raw.decode_content = True

    with lock:
        downloaded_images[str(idx)] = convert_webp_to_jpg(io.BytesIO(req.content))

    print(f"Downloaded {url}", f"using proxy {proxy}" if USE_PROXIES else "")


def convert_webp_to_jpg(webp_blob):
    webp_image = PIL.Image.open(webp_blob)
    rgb_image = webp_image.convert("RGB")
    jpeg_data = io.BytesIO()
    rgb_image.save(jpeg_data, format="JPEG")
    jpeg_blob = io.BytesIO(jpeg_data.getvalue())

    return jpeg_blob


def get_image_size(blob):
    img = PIL.Image.open(blob)
    width, height = img.size

    width /= 72
    height /= 72

    return width, height


def split_list_in_chunks(mylist: list, chunks_n: int):
    for i in range(0, len(mylist), chunks_n):
        yield mylist[i : i + chunks_n]


def get_all_slides_images(url):
    # url = "https://www.slideshare.net/gheorghio/la-seconda-guerra-mondiale-12879913"

    page_source = requests.get(url).text
    soup = bs4.BeautifulSoup(page_source, "html.parser")

    slides_number = scrape_slides_number(soup)
    ret = scrape_highest_img_src_set(soup)

    img_urls = []
    for i in range(1, slides_number + 1):
        a = ret.url.replace("SLIDE_NUMBER", str(i))
        img_urls.append(a)

    n = 4  # number of parallel connections
    chunks = list(split_list_in_chunks(img_urls, n))

    lock = threading.Lock()

    downloaded_images = dict()

    starting_timestamp = time.time()
    idx = 1
    for chunk in chunks:
        threads = []
        proxy = FreeProxy(rand=True).get()
        for c in chunk:
            thread = threading.Thread(
                target=save_image,
                args=(c, idx, downloaded_images, lock, proxy),
            )
            thread.start()
            threads.append(thread)

            idx += 1
        for thread in threads:
            thread.join()

    ending_timestamp = time.time()
    total_duration = int(ending_timestamp - starting_timestamp)
    print(f"Finished download using {n} threads in {total_duration} seconds")

    return downloaded_images, slides_number


app = FastAPI(docs_url=None, redoc_url=None, openapi_url=None)

origins = ["*"]
app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_methods=["GET"],
    allow_headers=["*"],
)

# get_all_slides_images("https://www.slideshare.net/Psicologi-A-Lavoro/slide-hikikomori")
# https://www.slideshare.net/slideshow/2doc-258758272/258758272 329 slides pptx


@app.get("/api/slidesharedl/ping")
async def pingpong():
    return "pong"


@app.get("/api/slidesharedl/download/{slideshare_url}")
async def download(slideshare_url: str, request: Request):
    try:
        client_ip = request.client.host
        decoded_url = requests.utils.unquote(decode_b64(slideshare_url))
        print(f"IP {client_ip} downloading {decoded_url}")

        downloaded_images, slides_number = get_all_slides_images(decoded_url)

        prs = pptx.Presentation()

        w, h = get_image_size(downloaded_images["1"])
        prs.slide_width = pptx.util.Inches(w)
        prs.slide_height = pptx.util.Inches(h)

        blank_slide_layout = prs.slide_layouts[6]

        left = top = pptx.util.Inches(0)
        width = prs.slide_width
        height = prs.slide_height

        for i in range(1, slides_number + 1):
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(
                downloaded_images[str(i)], left, top, width, height
            )

        presentation_bytes = io.BytesIO()
        prs.save(presentation_bytes)
        presentation_bytes.seek(0)

        # def iterfile():
        #    yield from presentation_bytes
        # return StreamingResponse(
        #    iterfile(),
        #    media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        #    headers={"Access-Control-Allow-Origin": "*"},
        # )

        return Response(
            content=presentation_bytes.read(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except:
        raise HTTPException(status_code=500, detail="Internal server error")


def get_slideshare_info(url):
    page_source = requests.get(url).text
    soup = bs4.BeautifulSoup(page_source, "html.parser")

    page_title = soup.find_all("title")[0].text.split(" | ")[0]
    img_src_set = scrape_highest_img_src_set(soup)
    slides_number = scrape_slides_number(soup)

    return {
        "title": page_title,
        "template_url": img_src_set.url,
        "slides_number": slides_number,
        "estimated_download_time_seconds": slides_number * 2,
    }


@app.get("/api/slidesharedl/info/{slideshare_url}")
async def info(slideshare_url: str, request: Request):
    try:
        client_ip = request.client.host
        decoded_url = requests.utils.unquote(decode_b64(slideshare_url))

        ret = get_slideshare_info(decoded_url)
        ret["client_ip"] = client_ip

        return ret
    except:
        raise HTTPException(status_code=500, detail="Internal server error")


# WIP
@app.get("/api/slidesharedl/get_slide/{slideshare_url}/{slide_number}")
async def get_slide(slideshare_url: str, slide_number: int, request: Request):
    try:
        decoded_url = requests.utils.unquote(decode_b64(slideshare_url))
        info = get_slideshare_info(decoded_url)
        template_url = info["template_url"]
        img_url = template_url.replace("SLIDE_NUMBER", str(slide_number))

        req = requests.get(img_url)

        return "👍"
    except:
        raise HTTPException(status_code=500, detail="Internal server error")


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=os.environ["PORT"] | 10000)
